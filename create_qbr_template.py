"""
QBR Template Generator for MSP Customer Success Tool (Business Impact Version)
Run this script to create Master_QBR_Template.pptx with the top 4 impact placeholders

Prerequisites:
    pip install python-pptx

Usage:
    python create_qbr_template.py
"""

import math

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define color scheme
BLUE = RGBColor(36, 46, 101)  # #242E65
GRAY = RGBColor(74, 85, 104)  # #4A5568
LIGHT_GRAY = RGBColor(226, 232, 240)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Quarterly Business Review"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Client Name
    client_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.7), Inches(8), Inches(0.8)
    )
    client_frame = client_box.text_frame
    client_frame.text = "{{CLIENT_NAME}}"
    client_frame.paragraphs[0].font.size = Pt(32)
    client_frame.paragraphs[0].font.color.rgb = GRAY
    client_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Review Period
    period_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    period_frame = period_box.text_frame
    period_frame.text = "{{REVIEW_PERIOD}}"
    period_frame.paragraphs[0].font.size = Pt(20)
    period_frame.paragraphs[0].font.color.rgb = GRAY
    period_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _estimate_text_height_in(text, font_pt, box_width_in, space_after_pt=0):
    """Estimate rendered height of a single text paragraph in inches.

    Uses average character width ≈ 0.55× font size (typical sans-serif).
    Errs slightly high (conservative), which prevents overlap.
    """
    chars_per_line = max(1, (box_width_in * 72) / (font_pt * 0.55))
    num_lines = max(1, math.ceil(len(text) / chars_per_line))
    line_height_pt = font_pt * 1.2  # standard line height multiplier
    return (num_lines * line_height_pt + space_after_pt) / 72


def add_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # --- Title ---
    TITLE_TOP = Inches(0.5)
    TITLE_HEIGHT = Inches(0.8)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), TITLE_TOP, Inches(9), TITLE_HEIGHT
    )
    title_frame = title_box.text_frame
    title_frame.text = "Executive Summary"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # --- Bullets ---
    BULLET_FONT_PT = 22
    BULLET_WIDTH_IN = 8
    SPACE_AFTER_PT = 14
    BULLETS_TOP = TITLE_TOP + TITLE_HEIGHT + Inches(0.2)

    bullets = [
        "{{SAME_DAY_RATE}}% of your team's IT issues were completely resolved on the same day.",
        "When employees reached out, we began working on their issues in an average of {{AVG_FIRST_RESPONSE}}.",
        "Critical business-halting emergencies were resolved in {{CRITICAL_RES_TIME}} on average.",
        "We managed {{TICKET_COUNT}} total IT events this quarter to keep your business running.",
    ]

    bullets_height_in = sum(
        _estimate_text_height_in(
            "• " + b, BULLET_FONT_PT, BULLET_WIDTH_IN, SPACE_AFTER_PT
        )
        for b in bullets
    )
    bullets_height_in = min(max(bullets_height_in, 2.5), 3.5)  # floor 2.5", cap 3.5"

    content_box = slide.shapes.add_textbox(
        Inches(1), BULLETS_TOP, Inches(BULLET_WIDTH_IN), Inches(bullets_height_in)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet_text
        p.font.size = Pt(BULLET_FONT_PT)
        p.font.color.rgb = GRAY
        p.space_after = Pt(SPACE_AFTER_PT)

    # --- Business Impact --- positioned below bullets
    IMPACT_TOP = BULLETS_TOP + Inches(bullets_height_in) + Inches(0.15)
    IMPACT_HEIGHT = Inches(0.4)
    impact_box = slide.shapes.add_textbox(
        Inches(0.5), IMPACT_TOP, Inches(9), IMPACT_HEIGHT
    )
    impact_frame = impact_box.text_frame
    impact_frame.word_wrap = True
    impact_frame.text = (
        "Business Impact: {{PRODUCTIVITY_HOURS_LOST}} productivity hours at risk"
        " | Est. cost: {{ESTIMATED_COST}}"
    )
    impact_frame.paragraphs[0].font.size = Pt(13)
    impact_frame.paragraphs[0].font.bold = True
    impact_frame.paragraphs[0].font.color.rgb = RGBColor(220, 38, 38)  # Red

    # --- Risk Statement --- positioned below Business Impact
    RISK_TOP = IMPACT_TOP + IMPACT_HEIGHT + Inches(0.1)
    RISK_HEIGHT = Inches(0.5)
    risk_box = slide.shapes.add_textbox(Inches(0.5), RISK_TOP, Inches(9), RISK_HEIGHT)
    risk_frame = risk_box.text_frame
    risk_frame.word_wrap = True
    risk_frame.text = "{{RISK_STATEMENT}}"
    risk_frame.paragraphs[0].font.size = Pt(12)
    risk_frame.paragraphs[0].font.italic = True
    risk_frame.paragraphs[0].font.color.rgb = GRAY

    # --- BEA box --- anchored below risk statement, never overlaps
    BEA_TOP = max(RISK_TOP + RISK_HEIGHT + Inches(0.15), Inches(6.3))
    LIGHT_BLUE = RGBColor(219, 234, 254)  # #DBEAFE
    bea_box = slide.shapes.add_shape(1, Inches(0.5), BEA_TOP, Inches(9), Inches(1.2))
    bea_box.fill.solid()
    bea_box.fill.fore_color.rgb = LIGHT_BLUE
    bea_box.line.color.rgb = BLUE

    # Line 1: Industry and GDP value
    line1_box = slide.shapes.add_textbox(
        Inches(0.65), BEA_TOP + Inches(0.05), Inches(8.7), Inches(0.45)
    )
    line1_frame = line1_box.text_frame
    line1_frame.word_wrap = True
    line1_frame.text = (
        "Industry Sector: {{BEA_INDUSTRY}}  |  "
        "GDP Value Added: {{BEA_LATEST_VALUE}} ({{BEA_LATEST_PERIOD}})"
    )
    line1_frame.paragraphs[0].font.size = Pt(13)
    line1_frame.paragraphs[0].font.bold = True
    line1_frame.paragraphs[0].font.color.rgb = BLUE

    # Line 2: Growth rates and trend label
    line2_box = slide.shapes.add_textbox(
        Inches(0.65), BEA_TOP + Inches(0.55), Inches(8.7), Inches(0.5)
    )
    line2_frame = line2_box.text_frame
    line2_frame.word_wrap = True
    line2_frame.text = (
        "QoQ Growth: {{BEA_QOQ_GROWTH}}  |  "
        "YoY Growth: {{BEA_YOY_GROWTH}}  |  {{BEA_TREND_LABEL}}"
    )
    line2_frame.paragraphs[0].font.size = Pt(12)
    line2_frame.paragraphs[0].font.color.rgb = GRAY


def add_metrics_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Key Business Impact Metrics"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Three metric cards for the new metrics
    metrics = [
        ("Same-Day Resolution", "{{SAME_DAY_RATE}}%"),
        ("Avg First Response", "{{AVG_FIRST_RESPONSE}}"),
        ("Critical Resolution", "{{CRITICAL_RES_TIME}}"),
    ]

    x_positions = [1, 3.7, 6.4]
    for i, (label, value) in enumerate(metrics):
        # Card background
        card = slide.shapes.add_shape(
            1, Inches(x_positions[i]), Inches(2.5), Inches(2.2), Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = BLUE

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x_positions[i]), Inches(2.7), Inches(2.2), Inches(0.8)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].font.size = Pt(14)
        label_frame.paragraphs[0].font.color.rgb = GRAY
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Value
        value_box = slide.shapes.add_textbox(
            Inches(x_positions[i]), Inches(3.5), Inches(2.2), Inches(1)
        )
        value_frame = value_box.text_frame
        value_frame.text = value
        value_frame.paragraphs[0].font.size = Pt(26)
        value_frame.paragraphs[0].font.bold = True
        value_frame.paragraphs[0].font.color.rgb = BLUE
        value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_chart_placeholder(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Support Type Distribution"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(9), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Proactive Maintenance vs Reactive Support"
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = GRAY

    # Placeholder text
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = "{{CHART_PLACEHOLDER}}"
    text_frame.paragraphs[0].font.size = Pt(24)
    text_frame.paragraphs[0].font.color.rgb = GRAY
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_stability_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "System Stability (Metric 1)"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Left box - Proactive
    left_box = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(3.5), Inches(2))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(34, 197, 94)  # Green

    left_text = slide.shapes.add_textbox(
        Inches(1), Inches(2.7), Inches(3.5), Inches(1.5)
    )
    left_frame = left_text.text_frame
    left_frame.text = "Proactive Work\n{{PROACTIVE_PERCENT}}%"
    left_frame.paragraphs[0].font.size = Pt(28)
    left_frame.paragraphs[0].font.bold = True
    left_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    left_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Right box - Reactive
    right_box = slide.shapes.add_shape(
        1, Inches(5.5), Inches(2.5), Inches(3.5), Inches(2)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(239, 68, 68)  # Red

    right_text = slide.shapes.add_textbox(
        Inches(5.5), Inches(2.7), Inches(3.5), Inches(1.5)
    )
    right_frame = right_text.text_frame
    right_frame.text = "Reactive Issues\n{{REACTIVE_PERCENT}}%"
    right_frame.paragraphs[0].font.size = Pt(28)
    right_frame.paragraphs[0].font.bold = True
    right_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    right_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Explanation
    exp_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    exp_frame = exp_box.text_frame
    exp_frame.text = "We actively prevent downtime before it impacts your employees. A higher proactive percentage means a more stable network."
    exp_frame.paragraphs[0].font.size = Pt(16)
    exp_frame.paragraphs[0].font.color.rgb = GRAY
    exp_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_responsiveness_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Responsiveness & Business Continuity"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Content bullets focused on Metrics 2, 3, and 4
    content_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(7), Inches(3.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    bullets = [
        "Speed to First Response: {{AVG_FIRST_RESPONSE}}",
        "Same-Day Resolution Rate: {{SAME_DAY_RATE}}%",
        "Critical Crisis Resolution Time: {{CRITICAL_RES_TIME}}",
    ]

    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet_text
        p.font.size = Pt(26)
        p.font.color.rgb = GRAY
        p.space_after = Pt(30)


def add_recommendations(prs, num_recommendations=3):
    """Builds the recommendations slide with dynamic slots."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Strategic Recommendations"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE

    # Risk Spotlight section
    RED = RGBColor(220, 38, 38)
    risk_header = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.0), Inches(9), Inches(0.3)
    )
    rh_frame = risk_header.text_frame
    rh_frame.text = "Risk Spotlight:"
    rh_frame.paragraphs[0].font.size = Pt(12)
    rh_frame.paragraphs[0].font.bold = True
    rh_frame.paragraphs[0].font.color.rgb = RED

    risk_text_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.35), Inches(9), Inches(0.75)
    )
    rt_frame = risk_text_box.text_frame
    rt_frame.word_wrap = True
    for i in range(3):
        p = rt_frame.paragraphs[0] if i == 0 else rt_frame.add_paragraph()
        p.text = f"{{{{TOP_RISK_{i + 1}}}}}"
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    # Dynamic vertical spacing based on number of recommendations
    usable_height = 5.0  # inches available below risk spotlight
    slot_height = min(usable_height / num_recommendations, 1.0)
    y_start = 2.2

    for i in range(num_recommendations):
        y_pos = y_start + (i * slot_height)

        # Number circle
        circle = slide.shapes.add_shape(
            1, Inches(0.5), Inches(y_pos), Inches(0.4), Inches(0.4)
        )
        circle.name = f"rec_{i + 1}_circle"
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE

        num_text = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos), Inches(0.4), Inches(0.4)
        )
        num_text.name = f"rec_{i + 1}_num"
        num_frame = num_text.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].font.size = Pt(16)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title placeholder
        title_ph = slide.shapes.add_textbox(
            Inches(1.1), Inches(y_pos), Inches(8.3), Inches(slot_height * 0.4)
        )
        title_ph.name = f"rec_{i + 1}_title"
        tf = title_ph.text_frame
        tf.text = f"{{{{REC_{i + 1}_TITLE}}}}"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLUE

        # Rationale placeholder
        rat_ph = slide.shapes.add_textbox(
            Inches(1.1),
            Inches(y_pos + slot_height * 0.42),
            Inches(8.3),
            Inches(slot_height * 0.5),
        )
        rat_ph.name = f"rec_{i + 1}_rationale"
        rf = rat_ph.text_frame
        rf.word_wrap = True
        rf.text = f"{{{{REC_{i + 1}_RATIONALE}}}}"
        rf.paragraphs[0].font.size = Pt(11)
        rf.paragraphs[0].font.color.rgb = GRAY


def add_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BLUE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Contact prompt
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.6))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Questions? Contact your account manager"
    contact_frame.paragraphs[0].font.size = Pt(22)
    contact_frame.paragraphs[0].font.color.rgb = GRAY
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Contact info placeholder
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = "{{MSP_CONTACT_INFO}}"
    info_frame.paragraphs[0].font.size = Pt(18)
    info_frame.paragraphs[0].font.color.rgb = GRAY
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Build all slides
    print("Creating Impact-focused slides...")
    add_title_slide(prs)
    add_executive_summary(prs)
    add_metrics_overview(prs)
    add_chart_placeholder(prs)
    add_stability_slide(prs)
    add_responsiveness_slide(prs)  # Replaced SLA performance slide with this
    add_recommendations(prs, num_recommendations=10)
    add_thank_you(prs)

    # Save the presentation
    filename = "Master_QBR_Template.pptx"
    prs.save(filename)
    print(f"✅ Business Impact QBR Template created successfully: {filename}")
    print("\n📋 Placeholders you need to populate with HaloPSA data:")
    placeholders = [
        "{{CLIENT_NAME}}",
        "{{REVIEW_PERIOD}}",
        "{{TICKET_COUNT}}",
        "{{PROACTIVE_PERCENT}}",
        "{{REACTIVE_PERCENT}}",
        "{{SAME_DAY_RATE}}",
        "{{AVG_FIRST_RESPONSE}}",
        "{{CRITICAL_RES_TIME}}",
        "{{CHART_PLACEHOLDER}}",
        "{{RECOMMENDATION_1}}",
        "{{RECOMMENDATION_2}}",
        "{{RECOMMENDATION_3}}",
        "{{MSP_CONTACT_INFO}}",
        "{{BEA_INDUSTRY}}",
        "{{BEA_LATEST_VALUE}}",
        "{{BEA_LATEST_PERIOD}}",
        "{{BEA_QOQ_GROWTH}}",
        "{{BEA_YOY_GROWTH}}",
        "{{BEA_TREND_LABEL}}",
        "{{PRODUCTIVITY_HOURS_LOST}}",
        "{{ESTIMATED_COST}}",
        "{{RISK_STATEMENT}}",
        "{{TOP_RISK_1}}",
        "{{TOP_RISK_2}}",
        "{{TOP_RISK_3}}",
    ]
    for p in placeholders:
        print(f"   • {p}")


if __name__ == "__main__":
    main()
