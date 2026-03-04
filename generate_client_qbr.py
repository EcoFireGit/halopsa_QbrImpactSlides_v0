"""
Client QBR Generator
Calculates business impact metrics from HaloPSA ticket data and populates the PPTX template.

Prerequisites:
    pip install python-pptx

Usage:
    python generate_client_qbr.py
"""

from datetime import datetime
from pptx import Presentation
import os

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import tempfile


def generate_support_distribution_chart(proactive_pct, reactive_pct, output_path):
    """
    Generates a horizontal bar chart showing proactive vs reactive support distribution.
    Saves it as a PNG to output_path. Returns the path for insertion into PPTX.
    """
    # Handle edge case where both are 0
    if proactive_pct == 0 and reactive_pct == 0:
        proactive_pct = 50
        reactive_pct = 50

    fig, ax = plt.subplots(figsize=(9, 3.5))
    fig.patch.set_facecolor("white")

    # Data
    categories = ["Support\nDistribution"]
    proactive_vals = [proactive_pct]
    reactive_vals = [reactive_pct]

    # Color scheme
    PROACTIVE_COLOR = "#22C55E"  # Green
    REACTIVE_COLOR = "#EF4444"  # Red
    BAR_HEIGHT = 0.5

    # Draw stacked horizontal bars
    ax.barh(
        categories,
        proactive_vals,
        height=BAR_HEIGHT,
        color=PROACTIVE_COLOR,
        label=f"Proactive ({int(proactive_pct)}%)",
    )
    ax.barh(
        categories,
        reactive_vals,
        height=BAR_HEIGHT,
        left=proactive_vals,
        color=REACTIVE_COLOR,
        label=f"Reactive ({int(reactive_pct)}%)",
    )

    # Add percentage labels inside the bars
    if proactive_pct >= 10:
        ax.text(
            proactive_pct / 2,
            0,
            f"{int(proactive_pct)}%",
            ha="center",
            va="center",
            fontsize=18,
            fontweight="bold",
            color="white",
        )

    if reactive_pct >= 10:
        ax.text(
            proactive_pct + reactive_pct / 2,
            0,
            f"{int(reactive_pct)}%",
            ha="center",
            va="center",
            fontsize=18,
            fontweight="bold",
            color="white",
        )

    # Style the chart
    ax.set_xlim(0, 100)
    ax.set_xlabel("Percentage of Total Tickets (%)", fontsize=11, color="#4A5568")
    ax.set_title(
        "Proactive vs. Reactive Support Distribution",
        fontsize=14,
        fontweight="bold",
        color="#242E65",
        pad=15,
    )

    # Remove chart borders
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)

    ax.tick_params(left=False)
    ax.xaxis.set_tick_params(labelsize=10, colors="#4A5568")
    ax.yaxis.set_ticklabels([])

    # Legend
    proactive_patch = mpatches.Patch(
        color=PROACTIVE_COLOR, label=f"Proactive ({int(proactive_pct)}%)"
    )
    reactive_patch = mpatches.Patch(
        color=REACTIVE_COLOR, label=f"Reactive ({int(reactive_pct)}%)"
    )
    ax.legend(
        handles=[proactive_patch, reactive_patch],
        loc="lower center",
        bbox_to_anchor=(0.5, -0.45),
        ncol=2,
        fontsize=11,
        frameon=False,
    )

    plt.tight_layout()
    plt.savefig(
        output_path, dpi=150, bbox_inches="tight", facecolor="white", edgecolor="none"
    )
    plt.close()
    print(f"✅ Chart saved to: {output_path}")
    return output_path


# --- 1. METRICS CALCULATION ENGINE ---


def calculate_metrics(tickets):
    """
    Computes the top 4 business impact metrics from raw ticket data.
    Hardened against all division-by-zero and data quality edge cases.
    """

    # --- Edge Case 1: None or non-list input ---
    if not tickets or not isinstance(tickets, list):
        print("⚠️  Warning: No ticket data provided. All metrics defaulted to N/A.")
        return _empty_metrics()

    total_tickets = len(tickets)

    # --- Edge Case 2: Empty list (returns complete dict, not partial) ---
    if total_tickets == 0:
        print("⚠️  Warning: Ticket list is empty. All metrics defaulted to N/A.")
        return _empty_metrics()

    # Configure based on your HaloPSA TicketType IDs
    PROACTIVE_TYPES = [30, 40, 100]
    REACTIVE_TYPES = [
        1,
        10,
        20,
        50,
        60,
        61,
        62,
        63,
        64,
        65,
        66,
        67,
        68,
        69,
        70,
        71,
        72,
        9999,
    ]

    proactive_count = 0
    reactive_count = 0
    closed_tickets = 0
    same_day_count = 0
    critical_tickets = 0
    critical_total_age = 0.0
    valid_response_tickets = 0
    total_response_minutes = 0.0

    for t in tickets:
        # --- Metric 1: Proactive vs Reactive ---
        tt_id = t.get("tickettype_id")
        if tt_id in PROACTIVE_TYPES:
            proactive_count += 1
        elif tt_id in REACTIVE_TYPES:
            reactive_count += 1

        # --- Metric 2: Same-Day Resolution ---
        # Edge Case 3: hasbeenclosed may not be a strict boolean
        # We use `is True` to avoid truthy strings like "true" or 1
        if t.get("hasbeenclosed") is True:
            closed_tickets += 1
            date_occurred = t.get("dateoccurred", "")
            date_closed = t.get("dateclosed", "")
            if date_occurred and date_closed:
                if date_occurred.split("T")[0] == date_closed.split("T")[0]:
                    same_day_count += 1

        # --- Metric 3: Critical Crisis Resolution ---
        if t.get("priority_id") == 1:
            raw_age = t.get("ticketage", 0.0)

            # Edge Case 4: ticketage can be negative (data sync issues in Halo)
            if isinstance(raw_age, (int, float)) and raw_age > 0:
                critical_tickets += 1
                critical_total_age += raw_age
            else:
                # Still count the critical ticket, but don't include it in time calc
                critical_tickets += 1

        # --- Metric 4: Speed to First Response ---
        date_occurred = t.get("dateoccurred", "")
        response_date = t.get("responsedate", "")

        if date_occurred and response_date and not response_date.startswith("0001"):
            try:
                t_occ = datetime.fromisoformat(date_occurred)
                t_res = datetime.fromisoformat(response_date)
                diff_minutes = (t_res - t_occ).total_seconds() / 60

                # Edge Case 5: response recorded BEFORE occurrence (clock skew)
                if diff_minutes >= 0:
                    total_response_minutes += diff_minutes
                    valid_response_tickets += 1
                else:
                    print(
                        f"⚠️  Skipping ticket {t.get('id')}: response date is before occurrence date."
                    )
            except ValueError as e:
                print(f"⚠️  Skipping ticket {t.get('id')}: invalid date format. ({e})")

    # --- Final Math (all denominators are now guaranteed non-zero) ---

    # Metric 1
    total_typed = proactive_count + reactive_count
    proactive_pct = (proactive_count / total_typed * 100) if total_typed > 0 else 0
    reactive_pct = (reactive_count / total_typed * 100) if total_typed > 0 else 0

    # Metric 2
    same_day_rate = (same_day_count / closed_tickets * 100) if closed_tickets > 0 else 0

    # Metric 3
    if critical_tickets > 0 and critical_total_age > 0:
        avg_crit_age = critical_total_age / critical_tickets
        crit_res_str = f"{avg_crit_age:.1f} hours"
    elif critical_tickets > 0 and critical_total_age == 0:
        # Critical tickets existed but all had invalid/negative ages
        crit_res_str = "< 1 hour"
    else:
        crit_res_str = "< 1 hour"

    # Metric 4
    if valid_response_tickets > 0:
        avg_resp_mins = total_response_minutes / valid_response_tickets
        if avg_resp_mins < 60:
            first_resp_str = f"{int(avg_resp_mins)} mins"
        else:
            first_resp_str = f"{avg_resp_mins / 60:.1f} hours"
    else:
        first_resp_str = "N/A"

    return {
        "{{TICKET_COUNT}}": str(total_tickets),
        "{{PROACTIVE_PERCENT}}": f"{int(proactive_pct)}",
        "{{REACTIVE_PERCENT}}": f"{int(reactive_pct)}",
        "{{SAME_DAY_RATE}}": f"{int(same_day_rate)}",
        "{{CRITICAL_RES_TIME}}": crit_res_str,
        "{{AVG_FIRST_RESPONSE}}": first_resp_str,
    }


def _empty_metrics():
    """Returns a complete placeholder dictionary with safe default values."""
    return {
        "{{TICKET_COUNT}}": "0",
        "{{PROACTIVE_PERCENT}}": "N/A",
        "{{REACTIVE_PERCENT}}": "N/A",
        "{{SAME_DAY_RATE}}": "N/A",
        "{{CRITICAL_RES_TIME}}": "N/A",
        "{{AVG_FIRST_RESPONSE}}": "N/A",
    }


def calculate_health_score(metrics_data: dict) -> int:
    """Compute a 0-100 Client Health Score from the calculate_metrics() dict."""

    # 1. Proactive % (0-25 pts) — higher is better
    try:
        score_proactive = (
            float(metrics_data.get("{{PROACTIVE_PERCENT}}", "0")) / 100.0
        ) * 25.0
    except (ValueError, TypeError):
        score_proactive = 0.0

    # 2. Same-day resolution rate (0-25 pts) — higher is better
    try:
        score_same_day = (
            float(metrics_data.get("{{SAME_DAY_RATE}}", "0")) / 100.0
        ) * 25.0
    except (ValueError, TypeError):
        score_same_day = 0.0

    # 3. Critical resolution time (0-25 pts) — lower is better
    # N/A or < 1 hour = 25; ≤4h = 25; 4–24h = linear; ≥24h = 0
    crit_str = metrics_data.get("{{CRITICAL_RES_TIME}}", "N/A").strip()
    if crit_str == "N/A" or crit_str.startswith("<"):
        score_critical = 25.0
    else:
        try:
            hours = float(crit_str.replace(" hours", "").replace(" hour", "").strip())
            if hours <= 4.0:
                score_critical = 25.0
            elif hours >= 24.0:
                score_critical = 0.0
            else:
                score_critical = 25.0 * (1.0 - (hours - 4.0) / 20.0)
        except (ValueError, TypeError):
            score_critical = 25.0

    # 4. Avg first response (0-25 pts) — lower is better
    # N/A = 12.5 (neutral); ≤30 min = 25; 30–240 min = linear; ≥240 min = 0
    resp_str = metrics_data.get("{{AVG_FIRST_RESPONSE}}", "N/A").strip()
    if resp_str == "N/A":
        score_response = 12.5
    else:
        try:
            if resp_str.endswith(" mins"):
                resp_minutes = float(resp_str.replace(" mins", "").strip())
            elif "hour" in resp_str:
                resp_minutes = (
                    float(resp_str.replace(" hours", "").replace(" hour", "").strip())
                    * 60.0
                )
            else:
                resp_minutes = None

            if resp_minutes is None:
                score_response = 12.5
            elif resp_minutes <= 30.0:
                score_response = 25.0
            elif resp_minutes >= 240.0:
                score_response = 0.0
            else:
                score_response = 25.0 * (1.0 - (resp_minutes - 30.0) / 210.0)
        except (ValueError, TypeError):
            score_response = 12.5

    return int(
        round(score_proactive + score_same_day + score_critical + score_response)
    )


# --- 2. TEMPLATE POPULATION ENGINE ---


def replace_text_in_shape(shape, replacements):
    """Recursively search for placeholders in PPTX shapes and replace them."""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, val in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, val)

    # Handle grouped shapes
    if shape.shape_type == 6:  # Group
        for child_shape in shape.shapes:
            replace_text_in_shape(child_shape, replacements)


def _remove_unused_rec_slots(slide, num_recs):
    """Delete the circle, number, title, and rationale shapes for slots N > num_recs."""
    if num_recs >= 10:
        return

    shapes_to_delete = []
    tops_to_delete = set()

    # First pass: find title/rationale textboxes for unused slots
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for n in range(num_recs + 1, 11):
                    if f"{{{{REC_{n}_TITLE}}}}" in run.text:
                        tops_to_delete.add(shape.top)
                        shapes_to_delete.append(shape)
                    elif f"{{{{REC_{n}_RATIONALE}}}}" in run.text:
                        shapes_to_delete.append(shape)

    # Second pass: delete circle + number textbox at same top as unused title shapes
    for shape in slide.shapes:
        if shape not in shapes_to_delete and shape.top in tops_to_delete:
            shapes_to_delete.append(shape)

    sp_tree = slide.shapes._spTree
    for shape in shapes_to_delete:
        sp_tree.remove(shape._element)


def generate_qbr(template_path, output_path, contextual_data, ticket_data, num_recs=10):
    """
    Loads the template, computes metrics, inserts the chart image,
    replaces all text placeholders, and saves the final PPTX.
    """

    if not os.path.exists(template_path):
        print(
            f"❌ Error: Cannot find {template_path}. Run create_qbr_template.py first."
        )
        return

    # 1. Compute metrics from tickets
    metrics_data = calculate_metrics(ticket_data)

    # 2. Generate the chart PNG to a temp file
    chart_path = tempfile.mktemp(suffix=".png")
    generate_support_distribution_chart(
        proactive_pct=float(metrics_data.get("{{PROACTIVE_PERCENT}}", 0) or 0),
        reactive_pct=float(metrics_data.get("{{REACTIVE_PERCENT}}", 0) or 0),
        output_path=chart_path,
    )

    # 3. Combine all text replacement data
    # Remove {{CHART_PLACEHOLDER}} from text replacements since we handle it visually
    final_replacements = {**contextual_data, **metrics_data}
    final_replacements.pop("{{CHART_PLACEHOLDER}}", None)

    # 4. Open the presentation
    prs = Presentation(template_path)

    # Remove shapes for unused recommendation slots
    for slide in prs.slides:
        _remove_unused_rec_slots(slide, num_recs)

    for slide in prs.slides:
        chart_inserted = False

        for shape in slide.shapes:
            # --- Check if this shape contains the chart placeholder text ---
            if shape.has_text_frame and not chart_inserted:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{CHART_PLACEHOLDER}}" in run.text:
                            # Get the position and size of the placeholder shape
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height

                            # Clear the placeholder text
                            run.text = ""

                            # Insert the chart image in the same position
                            slide.shapes.add_picture(
                                chart_path, left, top, width, height
                            )
                            chart_inserted = True
                            print(
                                f"✅ Chart inserted on slide: '{slide.shapes.title.text if slide.shapes.title else 'Untitled'}'"
                            )
                            break

            # --- Replace all other text placeholders ---
            replace_text_in_shape(shape, final_replacements)

    # 5. Save the final PPTX
    prs.save(output_path)
    print(f"✅ QBR saved to: {output_path}")

    # 6. Clean up the temp chart file
    if os.path.exists(chart_path):
        os.remove(chart_path)
        print("🧹 Temp chart file cleaned up.")


def build_recommendation_replacements(recommendations: list[dict]) -> dict:
    """
    Converts Claude's recommendation list into PPTX placeholder key-value pairs.
    Fills any unused slots (up to 10) with empty strings.
    """
    replacements = {}
    for i in range(10):  # Always generate keys for all 10 possible slots
        if i < len(recommendations):
            replacements[f"{{{{REC_{i + 1}_TITLE}}}}"] = recommendations[i]["title"]
            replacements[f"{{{{REC_{i + 1}_RATIONALE}}}}"] = recommendations[i][
                "rationale"
            ]
        else:
            # Clear unused slots so no placeholder text remains in the PPTX
            replacements[f"{{{{REC_{i + 1}_TITLE}}}}"] = ""
            replacements[f"{{{{REC_{i + 1}_RATIONALE}}}}"] = ""
    return replacements


# --- 3. EXECUTION WITH MOCK DATA ---

if __name__ == "__main__":
    # Mock data based on your exact HaloPSA JSON structure
    mock_tickets = [
        {  # Ticket 1: Reactive, Same-Day Resolution, 15 min response
            "id": 1001,
            "tickettype_id": 1,
            "priority_id": 3,
            "hasbeenclosed": True,
            "dateoccurred": "2026-02-17T09:00:00",
            "responsedate": "2026-02-17T09:15:00",
            "dateclosed": "2026-02-17T11:30:00",
            "ticketage": 2.5,
        },
        {  # Ticket 2: Proactive Alert
            "id": 1002,
            "tickettype_id": 4,
            "priority_id": 3,
            "hasbeenclosed": True,
            "dateoccurred": "2026-02-18T02:00:00",
            "responsedate": "2026-02-18T02:05:00",
            "dateclosed": "2026-02-18T03:00:00",
            "ticketage": 1.0,
        },
        {  # Ticket 3: Critical Crisis, Multi-day resolution
            "id": 1003,
            "tickettype_id": 1,
            "priority_id": 1,
            "hasbeenclosed": True,
            "dateoccurred": "2026-02-19T10:00:00",
            "responsedate": "2026-02-19T10:05:00",
            "dateclosed": "2026-02-21T14:00:00",
            "ticketage": 52.0,
        },
        {  # Ticket 4: Proactive Alert, Same-day
            "id": 1004,
            "tickettype_id": 4,
            "priority_id": 3,
            "hasbeenclosed": True,
            "dateoccurred": "2026-02-20T08:00:00",
            "responsedate": "2026-02-20T08:10:00",
            "dateclosed": "2026-02-20T09:00:00",
            "ticketage": 1.0,
        },
    ]

    # Compute business impact and risk flags from mock data
    from business_impact import (
        calculate_business_impact,
        format_impact_replacements,
    )
    from risk_analyzer import analyze_risks, format_risk_replacements
    from bea_insights import build_empty_bea_replacements

    mock_metrics = calculate_metrics(mock_tickets)
    impact = calculate_business_impact(mock_metrics, mock_tickets, 100, 65.0)
    risk_flags = analyze_risks(mock_tickets)

    mock_recs = [
        {
            "title": "Upgrade Legacy Workstations",
            "rationale": "Slow application performance causing productivity loss.",
        },
        {
            "title": "Security Awareness Training",
            "rationale": "Reduce phishing-related incidents for General Users.",
        },
        {
            "title": "Backup Redundancy Review",
            "rationale": "Critical servers need additional backup coverage.",
        },
    ]
    rec_replacements = build_recommendation_replacements(mock_recs)

    # Contextual data to fill the rest of the presentation
    context_data = {
        "{{CLIENT_NAME}}": "Acme Corporation",
        "{{REVIEW_PERIOD}}": "Q1 2026",
        "{{CHART_PLACEHOLDER}}": "[Chart Graphic Will Go Here]",
        "{{MSP_CONTACT_INFO}}": "Jane Doe | jdoe@yourmsp.com | (555) 123-4567",
        **rec_replacements,
        **build_empty_bea_replacements(),
        **format_impact_replacements(impact),
        **format_risk_replacements(risk_flags),
    }

    # Run the generator
    generate_qbr(
        template_path="Master_QBR_Template.pptx",
        output_path="Acme_Corp_Q1_2026_QBR.pptx",
        contextual_data=context_data,
        ticket_data=mock_tickets,
    )
