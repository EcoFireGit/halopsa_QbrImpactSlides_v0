"""
QBR Data Replacer - Replaces placeholders in Master_QBR_Template.pptx with real data

Prerequisites:
    pip install python-pptx

Usage:
    from qbr_data_replacer import replace_qbr_placeholders

    data = {
        "CLIENT_NAME": "Acme Corporation",
        "REVIEW_PERIOD": "Q1 2026 (January - March)",
        "TICKET_COUNT": "147",
        "EFFICIENCY_HOURS": "230",
        "AVG_RESOLUTION_TIME": "2.4",
        # ... more data
    }

    replace_qbr_placeholders("Master_QBR_Template.pptx", data, "Acme_Q1_2026_QBR.pptx")
"""

from pptx import Presentation
import os


def replace_text_in_shape(shape, replacements):
    """
    Recursively replace text in a shape and its sub-elements.
    Returns True if any replacement was made.
    """
    replaced = False

    # Check if shape has a text frame
    if hasattr(shape, "text_frame"):
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                new_text = original_text

                # Replace all placeholders in this run
                for placeholder, value in replacements.items():
                    placeholder_pattern = "{{" + placeholder + "}}"
                    if placeholder_pattern in new_text:
                        new_text = new_text.replace(placeholder_pattern, str(value))
                        replaced = True

                # Only update if text changed
                if new_text != original_text:
                    run.text = new_text

    # Check if shape has a table (for future enhancements)
    if hasattr(shape, "table"):
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        new_text = original_text

                        for placeholder, value in replacements.items():
                            placeholder_pattern = "{{" + placeholder + "}}"
                            if placeholder_pattern in new_text:
                                new_text = new_text.replace(
                                    placeholder_pattern, str(value)
                                )
                                replaced = True

                        if new_text != original_text:
                            run.text = new_text

    # Check for grouped shapes
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            if replace_text_in_shape(sub_shape, replacements):
                replaced = True

    return replaced


def replace_qbr_placeholders(template_path, data, output_path=None):
    """
    Replace placeholders in a PowerPoint template with real data.

    Args:
        template_path (str): Path to the template PPTX file
        data (dict): Dictionary mapping placeholder names to values
        output_path (str): Path for the output file. If None, overwrites template.

    Returns:
        str: Path to the generated file

    Example:
        data = {
            "CLIENT_NAME": "Acme Corp",
            "TICKET_COUNT": "147",
            "EFFICIENCY_HOURS": "230"
        }
        replace_qbr_placeholders("Master_QBR_Template.pptx", data, "Output.pptx")
    """

    # Validate template exists
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")

    # Load the presentation
    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)

    # Statistics
    total_replacements = 0
    slides_modified = 0

    # Iterate through all slides
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_had_replacement = False

        # Iterate through all shapes in the slide
        for shape in slide.shapes:
            if replace_text_in_shape(shape, data):
                slide_had_replacement = True
                total_replacements += 1

        if slide_had_replacement:
            slides_modified += 1
            print(f"  ✓ Slide {slide_num}: Replacements made")

    # Determine output path
    if output_path is None:
        output_path = template_path

    # Save the modified presentation
    prs.save(output_path)

    print("\n✅ Success!")
    print(f"   • Slides modified: {slides_modified}/{len(prs.slides)}")
    print(f"   • Total shape replacements: {total_replacements}")
    print(f"   • Output saved to: {output_path}")

    return output_path


def get_sample_data():
    """
    Returns sample data for testing the replacement function.
    This demonstrates the expected data structure.
    """
    return {
        "CLIENT_NAME": "Acme Corporation",
        "REVIEW_PERIOD": "Q1 2026 (January - March)",
        "TICKET_COUNT": "147",
        "EFFICIENCY_HOURS": "230",
        "AVG_RESOLUTION_TIME": "2.4",
        "CHART_PLACEHOLDER": "[Chart will be inserted here]",
        "PROACTIVE_PERCENT": "65",
        "REACTIVE_PERCENT": "35",
        "CRITICAL_COUNT": "12",
        "SLA_COMPLIANCE": "98.5",
        "RECOMMENDATION_1": "Consider upgrading legacy workstations to improve security posture",
        "RECOMMENDATION_2": "Implement monthly security awareness training for end users",
        "RECOMMENDATION_3": "Schedule quarterly network infrastructure health assessments",
        "MSP_CONTACT_INFO": "John Smith • jsmith@mspprovider.com • (555) 123-4567",
    }


def validate_data(data, required_placeholders=None):
    """
    Validates that all required placeholders have data.

    Args:
        data (dict): The data dictionary to validate
        required_placeholders (list): List of required placeholder names

    Returns:
        tuple: (is_valid, missing_fields)
    """
    if required_placeholders is None:
        required_placeholders = [
            "CLIENT_NAME",
            "REVIEW_PERIOD",
            "TICKET_COUNT",
            "EFFICIENCY_HOURS",
            "AVG_RESOLUTION_TIME",
            "PROACTIVE_PERCENT",
            "REACTIVE_PERCENT",
            "CRITICAL_COUNT",
            "SLA_COMPLIANCE",
            "RECOMMENDATION_1",
            "RECOMMENDATION_2",
            "RECOMMENDATION_3",
            "MSP_CONTACT_INFO",
        ]

    missing = [
        p
        for p in required_placeholders
        if p not in data or data[p] is None or data[p] == ""
    ]

    return len(missing) == 0, missing


# Main execution for testing
if __name__ == "__main__":
    print("=== QBR Data Replacer Test ===\n")

    # Test with sample data
    sample_data = get_sample_data()

    # Validate data
    is_valid, missing = validate_data(sample_data)
    if not is_valid:
        print(f"⚠️  Warning: Missing data for: {', '.join(missing)}")
    else:
        print("✓ All required fields present\n")

    # Attempt replacement
    template_file = "Master_QBR_Template_v1.pptx"
    output_file = "Test_QBR_Output_v1.pptx"

    try:
        replace_qbr_placeholders(template_file, sample_data, output_file)
        print(f"\n🎉 Test completed! Open {output_file} to review the results.")
    except FileNotFoundError as e:
        print(f"\n❌ Error: {e}")
        print(
            "\nMake sure you've run create_qbr_template.py first to generate the template."
        )
    except Exception as e:
        print(f"\n❌ Unexpected error: {e}")
